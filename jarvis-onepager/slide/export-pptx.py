#!/usr/bin/env python3
"""
slide/export-pptx.py — Optional PPTX export via python-pptx + Playwright screenshot.

Requirements:
  pip install python-pptx pillow
  npx playwright install chromium

Usage:
  python export-pptx.py
  python export-pptx.py --out /path/to/jarvis-slide.pptx

Output: jarvis-slide.pptx in the slide/ directory (or --out path).

Manual alternative (if python-pptx not installed):
  1. Open slide/index.html in Chrome at exactly 1920×1080
  2. File → Print → Destination: Save as PDF (no margins, Landscape)
  3. Import the PDF into Keynote → File → Export To → PowerPoint
"""
import sys
import subprocess
import importlib.util
import argparse
from pathlib import Path

SLIDE_PATH = Path(__file__).parent / "index.html"
DEFAULT_OUT = Path(__file__).parent / "jarvis-slide.pptx"

def main():
    parser = argparse.ArgumentParser(description="Export slide to PPTX")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()

    if importlib.util.find_spec("pptx") is None:
        print("python-pptx not installed.")
        print("Install with: pip install python-pptx pillow")
        print("\nManual export path:")
        print("  1. Open slide/index.html in Chrome at 1920×1080")
        print("  2. Print → Save as PDF (no margins)")
        print("  3. Import into Keynote → Export as PPTX")
        sys.exit(0)

    from pptx import Presentation
    from pptx.util import Emu

    # Render to PNG via Playwright
    png_path = Path("/tmp/jarvis-slide-export.png")
    print(f"Rendering slide to PNG via Playwright...")
    result = subprocess.run(
        [
            "npx", "playwright", "screenshot",
            "--viewport-size=1920,1080",
            f"file://{SLIDE_PATH.resolve()}",
            str(png_path),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        print(f"Playwright screenshot failed: {result.stderr}")
        print("Install: npx playwright install chromium")
        sys.exit(1)

    # Build PPTX (16:9, 10in × 5.625in at 96dpi)
    prs = Presentation()
    prs.slide_width  = Emu(9_144_000)   # 10 inches
    prs.slide_height = Emu(5_143_500)   # 5.625 inches

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(png_path), 0, 0, prs.slide_width, prs.slide_height
    )

    prs.save(str(args.out))
    print(f"Saved: {args.out}")

if __name__ == "__main__":
    main()
